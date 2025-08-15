from flask import Flask, request, jsonify
from flask_cors import CORS
from pymongo import MongoClient
from dotenv import load_dotenv
import os
import bcrypt
import jwt
import datetime

import google.generativeai as genai
import json
import re
import requests
from bs4 import BeautifulSoup
from googlesearch import search
from concurrent.futures import ThreadPoolExecutor, as_completed

from pptx import Presentation
import PyPDF2
import mammoth

import sys
print("Python executable path:", sys.executable)

# Load .env variables
load_dotenv()
MONGO_URI = os.getenv("MONGO_URI")
JWT_SECRET = os.getenv("JWT_SECRET")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Init Flask and Mongo
app = Flask(__name__)
CORS(app, resources={
    r"/*": {
        "origins": [
            "http://localhost:5173",
            "https://ves-pied.vercel.app"
        ]
    }
}, supports_credentials=True)
client = MongoClient(MONGO_URI)
db = client["your-db-name"]
users = db["users"]

ALLOWED_ORIGINS = [
    "http://localhost:5173",
    "https://ves-pied.vercel.app"
]

@app.after_request
def add_cors_headers(response):
    origin = request.headers.get("Origin")
    if origin in ALLOWED_ORIGINS:
        response.headers["Access-Control-Allow-Origin"] = origin
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    return response

@app.errorhandler(Exception)
def handle_error(e):
    response = jsonify({"error": str(e)})
    response.status_code = getattr(e, "code", 500)
    origin = request.headers.get("Origin")
    if origin in ALLOWED_ORIGINS:
        response.headers["Access-Control-Allow-Origin"] = origin
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    return response


# Gemini setup
genai.configure(api_key=GEMINI_API_KEY)

@app.route("/api/register", methods=["POST"])
def register():
    data = request.json
    username = data.get("username")
    password = data.get("password")

    if users.find_one({"username": username}):
        return jsonify({"message": "Username already exists"}), 400

    hashed = bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt())
    users.insert_one({"username": username, "password": hashed})

    return jsonify({"message": "User registered successfully"}), 201


@app.route("/api/login", methods=["POST"])
def login():
    data = request.json
    username = data.get("username")
    password = data.get("password")

    user = users.find_one({"username": username})
    if not user or not bcrypt.checkpw(password.encode("utf-8"), user["password"]):
        return jsonify({"message": "Invalid credentials"}), 401

    token = jwt.encode({
        "id": str(user["_id"]),
        "username": user["username"],
        "exp": datetime.datetime.utcnow() + datetime.timedelta(hours=1)
    }, JWT_SECRET, algorithm="HS256")

    return jsonify({"token": token})

def format_json_input(json_string: str) -> dict:
    try:
        json_data = json.loads(json_string.strip())
        return {k: str(v).strip() if isinstance(v, str) else v for k, v in json_data.items()}
    except Exception as e:
        raise ValueError(f"Error: {e}")


@app.route('/generate-experiment', methods=["GET"])
def generate_experiment():
    text = request.args.get("text")
    if not text:
        return jsonify({"error": "Missing text parameter"}), 400

    sys_instruct = f"""
    You are a helpful assistant that generates a science experiment in JSON format based on the given topic.
    Generate a JSON object describing a simple science experiment related to the topic '{text}'.
    The JSON must include the following keys:
    - "aim" (string): The objective of the experiment.
    - "introduction" (string): A brief introduction to the experiment, explaining its context or relevance.
    - "article" (string): A detailed description of the experiment, including materials needed and step-by-step instructions.
    Wrap the output in ```json``` markers and ensure it is valid JSON.
    """

    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        system_instruction=sys_instruct
    )

    try:
        response = model.generate_content(text)
        # Safely extract text from Gemini response
        gemini_text = getattr(response, "text", None)
        if not gemini_text and hasattr(response, "candidates"):
            gemini_text = response.candidates[0].content.parts[0].text

        if not gemini_text:
            return jsonify({"error": "No content returned from Gemini", "raw": "No response text"}), 400

        print("Raw Gemini response:", gemini_text)  # Log for debugging

        # Clean the response
        cleaned = re.sub(r"```json|```", "", gemini_text).strip()
        json_data = json.loads(cleaned)  # Parse JSON

        # Format the response to match frontend expectations
        formatted = {
            "aim": json_data.get("aim", ""),
            "introduction": json_data.get("introduction", ""),
            "article": json_data.get("article", "")
        }
        return jsonify(formatted)

    except json.JSONDecodeError as e:
        return jsonify({"error": f"JSON decode error: {str(e)}", "raw": gemini_text}), 400
    except Exception as e:
        return jsonify({"error": str(e), "raw": gemini_text if 'gemini_text' in locals() else "No response"}), 500
    

@app.route('/generate-mcq', methods=["GET"])
def generate_mcq():
    topic = request.args.get("topic", "general knowledge")

    sys_instruct = f"""
    You are a helpful assistant that creates multiple-choice questions (MCQs) in JSON format.
    Generate exactly 5 questions about the topic '{topic}'.
    Output a JSON array only, where each element has:
    "id" (integer), "question" (string), "options" (array of strings),
    "answer" (string, must match one of the options), and "difficulty" ("easy", "medium", or "hard").
    """

    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        system_instruction=sys_instruct
    )

    response = model.generate_content(topic)

    # Try to extract text from Gemini response safely
    try:
        gemini_text = getattr(response, "text", None)
        if not gemini_text and hasattr(response, "candidates"):
            gemini_text = response.candidates[0].content.parts[0].text

        if not gemini_text:
            return jsonify({"error": "No content from Gemini"}), 400

        cleaned = re.sub(r"```json|```", "", gemini_text).strip()
        data = json.loads(cleaned)

        # If Gemini returned just an array
        if isinstance(data, list):
            questions = data
        elif isinstance(data, dict) and "questions" in data:
            questions = data["questions"]
        else:
            return jsonify({"error": "Invalid JSON format", "raw": cleaned}), 400

        formatted = [{
            "id": q["id"],
            "question": q["question"],
            "options": q["options"],
            "correctAnswer": q["answer"],
            "difficulty": q["difficulty"]
        } for q in questions]

        return jsonify({"totalQuestions": len(formatted), "questions": formatted})

    except json.JSONDecodeError as e:
        return jsonify({"error": f"JSON decode error: {e}", "raw": gemini_text}), 400
    except Exception as e:
        return jsonify({"error": str(e)}), 400


def fetch_url(url):
    try:
        r = requests.get(url, timeout=5, headers={"User-Agent": "Mozilla/5.0"})
        soup = BeautifulSoup(r.text, 'html.parser')
        title = soup.title.string if soup.title else "No title"
        desc = soup.find("meta", attrs={"name": "description"})
        if desc and "content" in desc.attrs:
            description = desc["content"]
        else:
            description = soup.find("p").get_text() if soup.find("p") else "No description"
        return {"title": title, "link": url, "description": description}
    except:
        return None


@app.route('/search', methods=['GET'])
def search_api():
    query = request.args.get('query')
    if not query:
        return jsonify({"error": "Missing query"}), 400

    urls = list(search("learn " + query, num_results=10))
    results = []

    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = [executor.submit(fetch_url, url) for url in urls]
        for f in as_completed(futures):
            res = f.result()
            if res: results.append(res)

    return jsonify(results)

@app.route('/extract-text', methods=['POST'])
def extract_text():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    path = os.path.join("uploads", filename)
    os.makedirs("uploads", exist_ok=True)
    file.save(path)

    try:
        if ext == '.pptx':
            prs = Presentation(path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif ext == '.pdf':
            pdf = PyPDF2.PdfReader(open(path, 'rb'))
            text = "\n".join([p.extract_text() for p in pdf.pages])
        elif ext == '.docx':
            with open(path, 'rb') as docx_file:
                result = mammoth.convert_to_text(docx_file)
                text = result.value
        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400
        return jsonify({"filename": filename, "text": text})
    finally:
        os.remove(path)

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
